import datetime
import os

import pandas as pd

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "output")


def _ensure_output():
    os.makedirs(OUTPUT_DIR, exist_ok=True)


def create_ppt_from_data(dataframe: pd.DataFrame, output_path: str = None) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return "python-pptx not installed. Run: pip install python-pptx"

    if output_path is None:
        _ensure_output()
        ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = os.path.join(OUTPUT_DIR, f"report_{ts}.pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "FRIDAY Data Report"
    slide.placeholders[
        1
    ].text = f"Generated {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}"

    # Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Summary"
    rows, cols = dataframe.shape
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))
    tf = textbox.text_frame
    tf.text = f"Rows: {rows}, Columns: {cols}\n\nColumns:\n" + "\n".join(
        f"  • {col} ({dtype})"
        for col, dtype in zip(dataframe.columns, dataframe.dtypes.astype(str))
    )

    # Data table slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Data Preview"
    n_rows = min(rows, 20)
    n_cols = min(cols, 8)

    table = slide.shapes.add_table(
        n_rows + 1, n_cols, Inches(1), Inches(1.5), Inches(11), Inches(5)
    ).table

    for j, col in enumerate(dataframe.columns[:n_cols]):
        cell = table.cell(0, j)
        cell.text = str(col)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)

    for i in range(n_rows):
        for j in range(n_cols):
            cell = table.cell(i + 1, j)
            cell.text = str(dataframe.iloc[i, j])[:20]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)

    prs.save(output_path)
    return f"PPT generated -> {output_path}"


def add_slide_with_chart(title: str, chart_image: str) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return "python-pptx not installed"

    _ensure_output()
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = os.path.join(OUTPUT_DIR, f"chart_slide_{ts}.pptx")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    if os.path.isfile(chart_image):
        slide.shapes.add_picture(chart_image, Inches(1), Inches(1.5), width=Inches(8))

    prs.save(output_path)
    return f"Slide with chart -> {output_path}"
